from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image,ImageDraw
import os 
import glob

prs = Presentation()
# this image directory file getting single image taking help to glob

def water_mark():
	for orginal_img in glob.glob("C:/Users/revan/Desktop/PPT Home Assignment/*.jpg"):
		nm = orginal_img.split('\\')# ------> This Will spilt the image path stor into list why i spilted becouse<----# 
		nam = nm[1].split('.')
		name = nam[0]
		position=(50,150)
		base_image = Image.open(orginal_img)
		watermark = Image.open('nike_black.png')
		size =(1500,1200)
		watermark.thumbnail(size)
		base_image.paste(watermark,position,mask=watermark)

	# ------ its all watermarked images will saved the particuler file-----#
		img_path ="C:/Users/revan/Desktop/PPT Home Assignment/imges"
		base_image.save(os.path.join(img_path,str(name)+".jpg"))

	# ------Here Below code starts PPT------------------#
	def PPT_slide():
		for images in glob.glob("C:/Users/revan/Desktop/PPT Home Assignment/imges/*.jpg"):
			nm = images.split('\\')# This Will spilt the image path stor into list why i spilted becouse 
			nam = nm[1].split('.')
			name = nam[0]

			slide_layout = prs.slide_layouts[1]
			slide = prs.slides.add_slide(slide_layout)
			shapes = slide.shapes

			title_shape = shapes.title
			body_shape = shapes.placeholders[1]
			title_shape.text = 'This Images slides is '+str(name)
			tf = body_shape.text_frame
			tf.text = 'This The '+str(name) 
			right=right=Inches(0.5)
			top=Inches(2.5)
			height=Inches(4)
			width=Inches(4)
			img = shapes.add_picture(images,right,top,width,height)
		prs.save('PPT_Home_Assignment.pptx')
		os.startfile('PPT_Home_Assignment.pptx')
	return PPT_slide()
if __name__ == '__main__':
	water_mark()